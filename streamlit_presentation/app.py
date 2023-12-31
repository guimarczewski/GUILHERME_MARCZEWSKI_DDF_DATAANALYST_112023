import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import openai
from openai import OpenAI
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
from io import BytesIO
import re
import ast
import requests

@st.cache
def load_data(file_url):
    """functio to load csv file"""
    data = pd.read_csv(file_url, sep=';')
    # Substituir valores nulos na coluna 'category'
    data['category'].fillna("No department defined", inplace=True)
    return data

def replace_text(replacements, slide):
    """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text

def generate_gpt_response(gpt_input, max_tokens):
    """function to generate a response from GPT-3. Takes input and max tokens as arguments and returns a response"""
    # Create an instance of the OpenAI class
    chat = ChatOpenAI(openai_api_key=st.secrets["openai_credentials"]["API_KEY"], model='gpt-4-1106-preview',
                      temperature=0, max_tokens=max_tokens)

    # Generate a response from the model
    response = chat.predict_messages(
        [SystemMessage(content='You are a helpful expert in products and sales.'),
         HumanMessage(
             content=gpt_input)])

    return response.content.strip()

def generate_gpt_image(gpt_image_prompt, output_path="generated_image.jpg"):
    """Function to generate an image from GPT-3.5 and save it as a JPEG file."""

    client = OpenAI(api_key=st.secrets["openai_credentials"]["API_KEY"])

    # Create an instance of the OpenAI class
    response = client.images.generate(
        model="dall-e-2",  # Specify the image model
        prompt=gpt_image_prompt,
        n=1,
        size="512x512",
        response_format='url'
    )

    # Get the image URL from the response
    image_url = response.data[0].url

    # Download the image
    image_data = requests.get(image_url).content

    # Save the image as a JPEG file
    with open(output_path, 'wb') as f:
        f.write(image_data)

    return output_path

def add_image(slide, image, left, top, width):
    """function to add an image to the PowerPoint slide and specify its position and width"""
    slide.shapes.add_picture(image, left=left, top=top, width=width)

def dict_from_string(response):
    """function to parse GPT response with competitors tickers and convert it to a dict"""
    # Find a substring that starts with '{' and ends with '}', across multiple lines
    match = re.search(r'\{.*?\}', response, re.DOTALL)

    dictionary = None
    if match:
        try:
            # Try to convert substring to dict
            dictionary = ast.literal_eval(match.group())
        except (ValueError, SyntaxError):
            # Not a dictionary
            return None
    return dictionary

# Função principal do aplicativo
def app():
    file_url = "streamlit_presentation/product-search-corpus-final.csv"
    ppt_file = "streamlit_presentation/template.pptx"

    # Load data from CSV file
    data = load_data(file_url)

    st.subheader("Select the category and product.")
    category_filter = st.selectbox("Filter by Category", data['category'].unique())
    filtered_data = data[data['category'] == category_filter]
    selected_product = st.selectbox("Select a Product", filtered_data['title'].unique())

    if ppt_file:
        # Add button to generate new PowerPoint file
        generate_new_ppt_button = st.button("Generate New PowerPoint")

        if generate_new_ppt_button:

            # Generate product name for GPT
            gpt_input = f"{selected_product} - give me a resumed name for this product with a maximum of 50 characters, with no more comments, just the name."
            product_name = generate_gpt_response(gpt_input, 60)

            description = filtered_data[filtered_data['title'] == selected_product]['text'].iloc[0]
            features = filtered_data[filtered_data['title'] == selected_product]['features'].iloc[0]

            # Generate main features
            gpt_input = f"product:{selected_product}, description:{description} - Give me a summary of the main features of this product  with a maximum of 100 characters, with no more comments, just the features."
            main_features = generate_gpt_response(gpt_input, 100)

            # Generate target audience
            gpt_input = f"{selected_product} - give me a target audience for this product, following these rules: 'Gender: male, female or both. Age: min-max. give me just this, with no more comments."
            target_audience = generate_gpt_response(gpt_input, 60)

            category_value = filtered_data[filtered_data['title'] == selected_product]['category'].iloc[0]

            gpt_input = f"product:{selected_product}, category:{category_value} and description: {description} - Using funnel ads strategy, create 4 descriptions for online ads: 2 for cold leads called cold_1, cold_2, 1 for remarketing called remarketing_1 and 1 for customers who abandoned the cart, called abandon_1. The strategies names are the exact names I gave . I want you to return in that order. Return output as a Python dictionary with strategic name as key and description as value.Do not return anything else."
            strategies = generate_gpt_response(gpt_input, 300)

            strategies_dict = dict_from_string(strategies)

            # Generate image
            gpt_input = f"product:{selected_product}, category:{category_value} and description: {description} - With this information, I want a brief description of the main characteristics for generating an image with dall-e-2. Therefore, it cannot exceed 900 characters. respond with only the main characteristics to include in the prompt with this limit"
            image_prompt_openai = generate_gpt_response(gpt_input, 900)

            image_prompt = f"description:{image_prompt_openai} - Generate a image that represents the product, only the product image with no more details or anything else"
            image_path = generate_gpt_image(image_prompt, output_path="streamlit_presentation/generated_image.jpg")

            presentation = Presentation(ppt_file)

            # Generate new PowerPoint file
            slide_0 = presentation.slides[0]
            replace_text({"{product_name}": product_name}, slide_0)

            cold_1 = str(strategies_dict['cold_1'])
            cold_2 = str(strategies_dict['cold_2'])
            remarketing = str(strategies_dict['remarketing_1'])
            abandoned = str(strategies_dict['abandon_1'])

            slide_1 = presentation.slides[1]
            replace_text({"{c}": product_name}, slide_1)
            replace_text({"{s}": category_value}, slide_1)
            replace_text({"{i}": target_audience}, slide_1)
            replace_text({"{f}": main_features}, slide_1)
            add_image(slide_1, image_path, left=Inches(6.5), width=Inches(3), top=Inches(2))

            slide_2 = presentation.slides[2]
            replace_text({"{s}": cold_1}, slide_2)
            replace_text({"{w}": cold_2}, slide_2)
            replace_text({"{o}": remarketing}, slide_2)
            replace_text({"{t}": abandoned}, slide_2)

            # Save updated presentation to BytesIO
            updated_ppt_bytesio = BytesIO()
            presentation.save(updated_ppt_bytesio)
            updated_ppt_bytes = updated_ppt_bytesio.getvalue()

            # Add download button for updated file
            st.download_button(
                label="Download Updated PowerPoint",
                data=updated_ppt_bytes,
                file_name="updated_ppt.pptx",
                key='download_button'
            )
            st.success("PowerPoint file updated successfully!")

# Executar o aplicativo
if __name__ == "__main__":
    app()
